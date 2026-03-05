"""
File Content Extractor
───────────────────────
Extracts plain text from Office documents, PDFs, and text files.
Called during SharePoint sync when a file is downloaded.

Supported formats:
  .docx  — Word documents        (python-docx)
  .pptx  — PowerPoint files      (python-pptx)
  .xlsx  — Excel spreadsheets    (openpyxl)
  .pdf   — PDF documents         (pdfplumber)
  .txt   — Plain text
  .md    — Markdown
  .csv   — CSV files
  .html  — HTML (stripped)

Returns plain text string, or empty string if extraction fails.
"""

import io
import re
import logging

logger = logging.getLogger(__name__)

# Max characters to extract per file — keeps content searchable without
# storing enormous amounts in MongoDB
MAX_CHARS = 50_000


def extract_docx(content: bytes) -> str:
    """Extract text from Word document bytes."""
    try:
        from docx import Document
        doc = Document(io.BytesIO(content))
        parts = []

        # Paragraphs
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)

        # Tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    parts.append(row_text)

        return "\n".join(parts)[:MAX_CHARS]
    except Exception as e:
        logger.warning(f"docx extraction failed: {e}")
        return ""


def extract_pptx(content: bytes) -> str:
    """Extract text from PowerPoint file bytes."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_parts = []

            for shape in slide.shapes:
                # Text frames (title, body, text boxes)
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_parts.append(text)

                # Tables inside slides
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            slide_parts.append(row_text)

            # Slide notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_parts.append(f"[Notes: {notes}]")

            if slide_parts:
                parts.append(f"[Slide {slide_num}] " + " ".join(slide_parts))

        return "\n".join(parts)[:MAX_CHARS]
    except Exception as e:
        logger.warning(f"pptx extraction failed: {e}")
        return ""


def extract_xlsx(content: bytes) -> str:
    """Extract text from Excel spreadsheet bytes."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        parts = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_parts = [f"[Sheet: {sheet_name}]"]
            row_count = 0

            for row in ws.iter_rows(values_only=True):
                # Skip completely empty rows
                cell_values = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if cell_values:
                    sheet_parts.append(" | ".join(cell_values))
                    row_count += 1
                if row_count >= 500:   # cap rows per sheet
                    sheet_parts.append("... (truncated)")
                    break

            if row_count > 0:
                parts.append("\n".join(sheet_parts))

        wb.close()
        return "\n\n".join(parts)[:MAX_CHARS]
    except Exception as e:
        logger.warning(f"xlsx extraction failed: {e}")
        return ""


def extract_pdf(content: bytes) -> str:
    """Extract text from PDF bytes."""
    try:
        import pdfplumber
        parts = []
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    parts.append(text.strip())
                if page_num >= 100:    # cap at 100 pages
                    parts.append("... (truncated at 100 pages)")
                    break
        return "\n\n".join(parts)[:MAX_CHARS]
    except Exception as e:
        logger.warning(f"pdf extraction failed: {e}")
        return ""


def extract_html(content: bytes) -> str:
    """Strip HTML tags from HTML file bytes."""
    try:
        text = content.decode("utf-8", errors="ignore")
        text = re.sub(r"<script[^>]*>.*?</script>", " ", text, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r"<style[^>]*>.*?</style>",  " ", text, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r"<[^>]+>", " ", text)
        text = re.sub(r"&nbsp;", " ", text)
        text = re.sub(r"&amp;",  "&", text)
        text = re.sub(r"&lt;",   "<", text)
        text = re.sub(r"&gt;",   ">", text)
        text = re.sub(r"\s+",    " ", text).strip()
        return text[:MAX_CHARS]
    except Exception as e:
        logger.warning(f"html extraction failed: {e}")
        return ""


def extract_text(content: bytes, filename: str) -> str:
    """
    Route to the correct extractor based on file extension.
    Returns plain text or empty string.
    """
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""

    extractors = {
        "docx": extract_docx,
        "pptx": extract_pptx,
        "xlsx": extract_xlsx,
        "xls":  extract_xlsx,   # openpyxl can handle xls too
        "pdf":  extract_pdf,
        "html": extract_html,
        "htm":  extract_html,
    }

    if ext in extractors:
        return extractors[ext](content)

    # Plain text formats — just decode
    if ext in ("txt", "md", "csv", "log", "json", "xml", "yaml", "yml"):
        try:
            return content.decode("utf-8", errors="ignore")[:MAX_CHARS]
        except Exception:
            return ""

    # Unsupported
    return ""


# File types we attempt to index
SUPPORTED_EXTENSIONS = {
    "docx", "pptx", "xlsx", "xls",
    "pdf",
    "txt", "md", "csv", "log",
    "html", "htm",
}

MAX_FILE_SIZE_MB = 20   # skip files larger than this
